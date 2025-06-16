# text_extractor.py
import os
import tempfile
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup, XMLParsedAsHTMLWarning
import csv
# import pandas as pd  # not used for CSV/XLSX now
from ebooklib import epub
from odf import text, teletype
from odf.opendocument import load as load_odf
import extract_msg
from PIL import Image, UnidentifiedImageError
import pytesseract
import zipfile
import warnings

# Suppress XMLParsedAsHTMLWarning globally
warnings.filterwarnings("ignore", category=XMLParsedAsHTMLWarning)

# For XLSX streaming
try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

import sys, contextlib
import logging

logger = logging.getLogger(__name__)

# In text_extractor.py
import xlrd
import time

def extract_text_from_xls(path):
    """
    Extract text from old-style Excel .xls files via xlrd.
    Yields (label, text) where label indicates sheet name.
    Reads rows in streaming fashion, grouping rows into chunks to avoid
    building a huge string at once.
    """
    # Optional: start time if you want per-file timeouts
    start_time = time.time()
    try:
        # on_demand=True may reduce memory footprint for large files
        wb = xlrd.open_workbook(path, on_demand=True, formatting_info=False)
    except Exception as e:
        logger.warning(f"Failed to open .xls file {path} with xlrd: {e}")
        return

    try:
        # Iterate sheets
        for sheet in wb.sheets():
            sheet_name = sheet.name or "Sheet"
            # Stream rows in chunks, e.g., 500 rows per chunk
            chunk_size = 500
            buffer = []
            rows_processed = 0
            nrows = sheet.nrows
            for row_idx in range(nrows):
                # Optional time check: abort if too slow
                if time.time() - start_time > 60:  # e.g. 60s timeout
                    logger.warning(f"Timeout reading .xls {path} sheet {sheet_name}")
                    # Yield whatever is in buffer so far, then return
                    if buffer:
                        text_chunk = "\n".join(buffer)
                        if text_chunk.strip():
                            yield f"XLS:{sheet_name}", text_chunk
                    return

                try:
                    row = sheet.row(row_idx)
                except Exception as e:
                    logger.debug(f"Skipping row {row_idx} in {path} sheet {sheet_name}: {e}")
                    continue

                cells = []
                for cell in row:
                    try:
                        val = cell.value
                        # cell.ctype can be used to format types, but str(val) is usually okay
                        if val is not None and str(val).strip() != "":
                            cells.append(str(val))
                    except Exception:
                        continue
                if cells:
                    buffer.append(" ".join(cells))
                rows_processed += 1

                if rows_processed >= chunk_size:
                    text_chunk = "\n".join(buffer)
                    buffer = []
                    rows_processed = 0
                    if text_chunk.strip():
                        yield f"XLS:{sheet_name}", text_chunk

            # After all rows, yield leftover
            if buffer:
                text_chunk = "\n".join(buffer)
                if text_chunk.strip():
                    yield f"XLS:{sheet_name}", text_chunk

    except Exception as e:
        logger.warning(f"Error iterating .xls {path}: {e}", exc_info=True)
    finally:
        try:
            # xlrd’s Book object doesn't have a close method in older versions,
            # but on_demand=True releases cached data when sheets are closed.
            # If available, you can call wb.release_resources():
            if hasattr(wb, "release_resources"):
                wb.release_resources()
        except Exception:
            pass

@contextlib.contextmanager
def suppress_stderr():
    devnull = open(os.devnull, 'w')
    old_stderr = sys.stderr
    try:
        sys.stderr = devnull
        yield
    finally:
        sys.stderr = old_stderr
        devnull.close()

def extract_text_from_pdf(path):
    """
    Extract text from PDF pages, suppressing MuPDF stderr output and skipping corrupt pages.
    """
    doc = None
    try:
        with suppress_stderr():
            doc = fitz.open(path)
    except Exception as e:
        logger.warning(f"Cannot open PDF {path}: {e}")
        return

    try:
        for page_number in range(doc.page_count):
            try:
                with suppress_stderr():
                    page = doc.load_page(page_number)
                    text = page.get_text()
                if text and text.strip():
                    yield f"PDF page {page_number+1}", text
            except Exception as e:
                logger.debug(f"Skipping corrupt PDF page {page_number+1} in {path}: {e}")
                continue
    finally:
        try:
            doc.close()
        except Exception:
            pass

def extract_text_from_docx(path):
    """
    Extract text from DOCX: paragraphs and table cells. Skip if cannot open.
    """
    try:
        doc = Document(path)
    except Exception:
        return
    parts = []
    try:
        # Paragraphs
        for p in doc.paragraphs:
            try:
                if p.text and p.text.strip():
                    parts.append(p.text)
            except Exception:
                continue
        # Tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    try:
                        if cell.text and cell.text.strip():
                            parts.append(cell.text)
                    except Exception:
                        continue
    except Exception:
        # In case iterating paragraphs/tables fails
        pass
    text_all = "\n".join(parts)
    if text_all.strip():
        yield "DOCX", text_all

def extract_text_from_pptx(path):
    """
    Extract text from PPTX slides. Skip if cannot open.
    """
    try:
        prs = Presentation(path)
    except Exception:
        return
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    texts.append(shape.text)
            except Exception:
                continue
        combined = "\n".join(texts)
        if combined.strip():
            yield f"PPTX slide {i}", combined

def extract_text_from_html(path):
    try:
        with open(path, encoding="utf-8", errors="ignore") as f:
            content = f.read()
    except Exception:
        return
    try:
        # Use HTML parser. If many XML-like warnings occurred, switch to "html.parser"
        soup = BeautifulSoup(content, "lxml")
        text_all = soup.get_text(separator="\n")
        if text_all and text_all.strip():
            yield "HTML", text_all
    except Exception:
        # Optionally try fallback parser
        try:
            soup = BeautifulSoup(content, "html.parser")
            text_all = soup.get_text(separator="\n")
            if text_all and text_all.strip():
                yield "HTML", text_all
        except Exception:
            return
        
def extract_text_from_xml(path):
    """
    For .xml or XHTML files: parse as XML to avoid HTMLParser warnings.
    """
    try:
        with open(path, encoding="utf-8", errors="ignore") as f:
            content = f.read()
    except Exception:
        return
    try:
        soup = BeautifulSoup(content, "lxml-xml")
        text_all = soup.get_text(separator="\n")
        if text_all and text_all.strip():
            yield "XML", text_all
    except Exception:
        return

def extract_text_from_xlsx(path):
    """
    Stream rows; caller checks elapsed after each chunk.
    """
    if load_workbook is None:
        return
    try:
        # Suppress openpyxl warnings during load
        with warnings.catch_warnings():
            warnings.simplefilter("ignore", category=UserWarning)
            wb = load_workbook(filename=path, read_only=True, data_only=True)
    except Exception as e:
        logger.warning(f"Failed to open XLSX {path}: {e}")
        return

    try:
        for sheet in wb.worksheets:
            # Collect rows in chunks of e.g. 500 rows
            buffer = []
            rows_read = 0
            for row in sheet.iter_rows(values_only=True):
                try:
                    cells = [str(cell) for cell in row if cell is not None]
                except Exception:
                    cells = []
                    for cell in row:
                        try:
                            cells.append(str(cell))
                        except Exception:
                            continue
                if cells:
                    buffer.append(" ".join(cells))
                rows_read += 1
                if rows_read >= 500:
                    text_chunk = "\n".join(buffer)
                    buffer = []
                    rows_read = 0
                    if text_chunk.strip():
                        yield f"XLSX:{sheet.title}", text_chunk
            # leftover
            if buffer:
                text_chunk = "\n".join(buffer)
                if text_chunk.strip():
                    yield f"XLSX:{sheet.title}", text_chunk
    except Exception as e:
        logger.warning(f"Error iterating XLSX {path}: {e}")
    finally:
        try:
            wb.close()
        except Exception:
            pass

def extract_text_from_csv(path):
    """
    Stream rows for CSV: read via csv.reader line-by-line.
    Yield chunks of N rows joined.
    """
    chunk_size = 1000  # Process 1000 rows at a time
    try:
        with open(path, newline='', encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            buffer = []
            for row in reader:
                buffer.append(", ".join(row))
                if len(buffer) >= chunk_size:
                    text_chunk = "\n".join(buffer)
                    buffer = []
                    if text_chunk.strip():
                        yield "CSV_chunk", text_chunk
            # Yield any remaining rows in the buffer
            if buffer:
                text_chunk = "\n".join(buffer)
                if text_chunk.strip():
                    yield "CSV_chunk", text_chunk
    except Exception as e:
        logger.warning(f"Error reading CSV {path}: {e}")
        # Fallback: read raw if streaming fails
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                if content.strip():
                    yield "CSV_raw", content
        except Exception as fallback_e:
            logger.error(f"Fallback failed for CSV {path}: {fallback_e}")

def extract_text_from_txt(path):
    """
    Stream lines for large text files: yield in chunks of N lines or bytes.
    This extractor yields segments labeled “TXT_part” for detection.
    """
    try:
        # Read in chunks of lines, e.g., 1000 lines at a time
        with open(path, encoding="utf-8", errors="ignore") as f:
            buffer = []
            lines_read = 0
            for line in f:
                buffer.append(line.rstrip("\n"))
                lines_read += 1
                if lines_read >= 1000:
                    text_chunk = "\n".join(buffer)
                    buffer = []
                    lines_read = 0
                    if text_chunk.strip():
                        yield "TXT_chunk", text_chunk
                    # Continue next chunk
            # leftover
            if buffer:
                text_chunk = "\n".join(buffer)
                if text_chunk.strip():
                    yield "TXT_chunk", text_chunk
    except Exception as e:
        # Fallback: read raw if streaming fails
        try:
            with open(path, "rb") as f2:
                data = f2.read(65536)
            text = data.decode("utf-8", errors="ignore")
            if text.strip():
                yield "TXT_fallback", text
        except Exception:
            return


def extract_text_from_epub(path):
    try:
        book = epub.read_epub(path)
    except Exception:
        return
    parts = []
    for item in book.get_items():
        if item.get_type() == epub.EpubHtml:
            try:
                content = item.get_content()
                # Parse as XML/XHTML
                soup = BeautifulSoup(content, "lxml-xml")
                text_part = soup.get_text(separator="\n")
                if text_part and text_part.strip():
                    parts.append(text_part)
            except Exception:
                continue
    text_all = "\n".join(parts)
    if text_all and text_all.strip():
        yield "EPUB", text_all

def extract_text_from_odt(path):
    """
    Extract from ODT via odfpy. Skip on errors.
    """
    try:
        odt = load_odf(path)
    except Exception:
        return
    parts = []
    try:
        for el in odt.getElementsByType(text.P):
            try:
                t = teletype.extractText(el)
                if t and t.strip():
                    parts.append(t)
            except Exception:
                continue
    except Exception:
        pass
    text_all = "\n".join(parts)
    if text_all.strip():
        yield "ODT", text_all

def extract_text_from_msg(path):
    """
    Extract from Outlook .msg files. Skip on errors.
    """
    try:
        msg = extract_msg.Message(path)
    except Exception:
        return
    try:
        if msg.subject:
            yield "MSG subject", msg.subject
    except Exception:
        pass
    try:
        if msg.body:
            yield "MSG body", msg.body
    except Exception:
        pass

def extract_text_from_image(path):
    """
    OCR via pytesseract. Skip if PIL cannot open or tesseract fails.
    """
    try:
        img = Image.open(path)
    except (UnidentifiedImageError, Exception):
        return
    try:
        text = pytesseract.image_to_string(img)
        if text and text.strip():
            yield "OCR image", text
    except Exception:
        return
    finally:
        try:
            img.close()
        except Exception:
            pass

def extract_text_from_zip(path):
    """
    Extract entries from ZIP/JAR (and ODF packages). Recursively apply extractors.
    Skip corrupted zips.
    """
    from config import EXTRACTORS  # avoid circular import
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            try:
                with zipfile.ZipFile(path, "r") as z:
                    for member in z.namelist():
                        if member.endswith("/"):
                            continue
                        try:
                            extracted_path = z.extract(member, tmpdir)
                        except Exception:
                            continue
                        ext = os.path.splitext(member)[1].lower()
                        extractor = EXTRACTORS.get(ext)
                        if extractor:
                            try:
                                for label, text in extractor(extracted_path):
                                    yield f"ZIP:{member}:{label}", text
                            except Exception:
                                continue
            except zipfile.BadZipFile:
                return
    except Exception:
        return

def fallback_raw_text_extraction(path):
    """
    Fallback method to extract raw text from any file type.
    """
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
            if content.strip():
                yield "RAW", content
    except Exception as e:
        logger.error(f"Fallback raw text extraction failed for {path}: {e}")
